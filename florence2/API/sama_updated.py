import os
import io
from docx import Document
from docx.oxml.ns import qn
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from PIL import Image
import pytesseract
import pymupdf
import pdfplumber  # For table extraction from PDFs
from pptx import Presentation
from io import StringIO
import configparser

config = configparser.ConfigParser()
config.read('config.properties')

# Load pytesseract path from config
pytesseract_path = config['pytesseract']['file_path']
pytesseract.pytesseract.tesseract_cmd = pytesseract_path

# Function to iterate through DOCX elements
def iter_block_items(parent):
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)

def extract_text_from_pptx(file_path):
    pptx_text = StringIO()
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pptx_text.write(shape.text)
    pptx_text.seek(0)
    return pptx_text.read()

# Function to read DOCX files
def read_docx(file_path):
    doc = Document(file_path)
    full_content = []

    # Define namespaces
    namespaces = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
        'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    }

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            paragraph_content = {
                'type': 'paragraph',
                'text': block.text,
                'style': block.style.name if block.style else None,
                'bold': False,
                'italic': False,
                'underline': False,
                'font_size': None,
                'font_color': None
            }

            # Check for text formatting
            for run in block.runs:
                if run.bold:
                    paragraph_content['bold'] = True
                if run.italic:
                    paragraph_content['italic'] = True
                if run.underline:
                    paragraph_content['underline'] = True
                if run.font.size:
                    paragraph_content['font_size'] = run.font.size
                if run.font.color and run.font.color.rgb:
                    paragraph_content['font_color'] = run.font.color.rgb

            full_content.append(paragraph_content)

            # Process inline shapes (which include images)
            for run in block.runs:
                drawing_elements = run._element.findall('.//w:drawing', namespaces=namespaces)
                for drawing in drawing_elements:
                    inline = drawing.find('.//wp:inline', namespaces=namespaces)
                    if inline is not None:
                        blip = inline.find('.//a:blip', namespaces=namespaces)
                        if blip is not None:
                            image_rid = blip.get(qn('r:embed'))
                            image_part = doc.part.related_parts[image_rid]
                            image_stream = io.BytesIO(image_part.blob)
                            image = Image.open(image_stream)

                            # Perform OCR on the image
                            text = pytesseract.image_to_string(image)

                            extent = inline.find('.//wp:extent', namespaces=namespaces)
                            width = int(extent.get('cx')) if extent is not None else 0
                            height = int(extent.get('cy')) if extent is not None else 0

                            full_content.append({
                                'type': 'image',
                                'text': text,
                                'width': width,
                                'height': height
                            })

        elif isinstance(block, Table):
            table_data = []
            for row in block.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            full_content.append({
                'type': 'table',
                'data': table_data
            })

    return full_content

# Function to process PDF files
def process_pdf(file_path):
    pdf_content = []

    # Open the PDF file
    with pymupdf.open(file_path) as pdf_document:
        for page_number in range(len(pdf_document)):
            page = pdf_document[page_number]
            text = page.get_text("text")
            pdf_content.append({
                'type': 'paragraph',
                'text': text,
                'page_number': page_number + 1
            })

            # Extract images from the PDF
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                image_stream = io.BytesIO(image_bytes)
                image = Image.open(image_stream)

                # OCR the image
                ocr_text = pytesseract.image_to_string(image)

                pdf_content.append({
                    'type': 'image',
                    'text': ocr_text,
                    'page_number': page_number + 1
                })

    # Process tables separately
    with pdfplumber.open(file_path) as pdf:
        for page_number, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            for table in tables:
                pdf_content.append({
                    'type': 'table',
                    'data': table,
                    'page_number': page_number + 1
                })

    return pdf_content

# Function to convert content to markdown
def convert_to_markdown(content):
    if content['type'] == 'paragraph':
        text = content.get('text', '')  # Get text safely, default to empty if missing

        # Only try to access 'style' if it's present
        style = content.get('style', None)

        if style and style.startswith('Heading'):
            try:
                level = int(''.join(filter(str.isdigit, style)))
            except ValueError:
                level = 1
            return f"{'#' * level} {text}\n\n"

        # Formatting options
        formatting = []
        if content.get('bold'):
            formatting.append('**')
        if content.get('italic'):
            formatting.append('*')
        if content.get('underline'):
            formatting.append('__')

        formatted_text = f"{''.join(formatting)}{text}{''.join(reversed(formatting))}"

        color_info = f"[Color: {content.get('font_color')}]" if content.get('font_color') else ""
        size_info = f"[Size: {content.get('font_size')}]" if content.get('font_size') else ""

        return f"{formatted_text} {color_info}{size_info}\n\n"

    elif content['type'] == 'table':
        table_md = ""
        for i, row in enumerate(content['data']):
            table_md += "| " + " | ".join(row) + " |\n"
            if i == 0:
                table_md += "|" + "|".join(["---" for _ in row]) + "|\n"
        return table_md + "\n"

    elif content['type'] == 'image':
        return f"[Image Content: {content.get('text', '')}]\n\n"

    # Fallback in case content type is not recognized
    return f"Unsupported content type: {content['type']}\n\n"

# Main execution
if __name__ == "__main__":
    # Specify the paths to your files (both DOCX and PDF)
    file_paths = ["testddoc.docx", "testddoc1.pdf"]

    for file_path in file_paths:
        print(f"Processing {file_path}")

        if file_path.endswith(".docx"):
            content = read_docx(file_path)
        elif file_path.endswith(".pdf"):
            content = process_pdf(file_path)
        elif file_path.endswith(".pptx"):
            content = extract_text_from_pptx(file_path)
        else:
            print(f"Unsupported file format: {file_path}")
            continue

        # Get the base filename without extension
        base_filename = os.path.splitext(os.path.basename(file_path))[0]

        # Create the output filename
        output_file = f"{base_filename}_converted.md"

        # Convert to markdown and write to a file
        with open(output_file, "w", encoding="utf-8") as f:
            for item in content:
                f.write(convert_to_markdown(item))

        print(f"Converted content has been written to '{output_file}'\n")

    print("All files have been processed and converted to markdown.")
